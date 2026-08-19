"""Generate the Akkaverse pitch deck (10 slides) as a polished, branded .pptx.

Run:  python docs/deck/build_deck.py
Output: docs/deck/Akkaverse-Pitch-Deck.pptx

Drop a real logo at docs/deck/assets/akka-logo.png to use it instead of the
drawn emblem. A family hero image is expected at docs/deck/assets/family-hero.png.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

# ---------------------------------------------------------------- brand palette
INK = RGBColor(0x0C, 0x0A, 0x14)        # near-black ink background
INK_2 = RGBColor(0x1A, 0x14, 0x2C)      # panel purple
SAFFRON = RGBColor(0xFF, 0x9F, 0x1C)    # kunkuma saffron (primary accent)
MAGENTA = RGBColor(0xE0, 0x3E, 0x8A)    # festive magenta
TEAL = RGBColor(0x2E, 0xC4, 0xB6)       # peacock teal
CREAM = RGBColor(0xF7, 0xF2, 0xE7)      # cream text
MUTED = RGBColor(0xB9, 0xB0, 0xC8)      # muted lavender text
CARD = RGBColor(0x24, 0x1C, 0x3A)       # card fill
LOGO_BLUE = RGBColor(0x1B, 0x4F, 0xB0)  # emblem ring blue
LOGO_RED = RGBColor(0xC8, 0x27, 0x2B)   # emblem Kannada red

FONT_HEAD = "Georgia"
FONT_BODY = "Segoe UI"

ASSETS = Path(__file__).with_name("assets")
LOGO = ASSETS / "akka-logo.png"
HERO = ASSETS / "family-hero.png"

EMU = 914400
W = int(13.333 * EMU)
H = int(7.5 * EMU)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

IN = lambda v: Emu(int(v * EMU))


def slide():
    return prs.slides.add_slide(BLANK)


def _no_line(shape):
    shape.line.fill.background()


def _flat(shape):
    shape.shadow.inherit = False


def bg(s, color=INK):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    r.fill.solid()
    r.fill.fore_color.rgb = color
    _no_line(r)
    _flat(r)
    s.shapes._spTree.remove(r._element)
    s.shapes._spTree.insert(2, r._element)
    return r


def band(s, x, y, w, h, color, alpha=None):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    r.fill.solid()
    r.fill.fore_color.rgb = color
    _no_line(r)
    _flat(r)
    if alpha is not None:
        _set_alpha(r, alpha)
    return r


def _set_alpha(shape, pct):
    sp = shape.fill._xPr.find(qn("a:solidFill"))
    srgb = sp.find(qn("a:srgbClr"))
    a = srgb.makeelement(qn("a:alpha"), {"val": str(int(pct * 1000))})
    srgb.append(a)


def glow(s, x, y, d, color, alpha=14):
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, IN(d), IN(d))
    c.fill.solid()
    c.fill.fore_color.rgb = color
    _no_line(c)
    _flat(c)
    _set_alpha(c, alpha)
    return c


def rounded(s, x, y, w, h, color, radius=0.08):
    r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    r.adjustments[0] = radius
    r.fill.solid()
    r.fill.fore_color.rgb = color
    _no_line(r)
    _flat(r)
    return r


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.0):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, item in enumerate(runs):
        txt, size, color, bold, font, italic = (list(item) + [False, FONT_BODY, False])[:6]
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.name = font
        r.font.color.rgb = color
    return tb


def dot(s, x, y, color, d=0.12):
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, IN(d), IN(d))
    c.fill.solid()
    c.fill.fore_color.rgb = color
    _no_line(c)
    _flat(c)
    return c


def ring(s, cx, cy, d, color, weight=2.0):
    o = s.shapes.add_shape(MSO_SHAPE.OVAL, IN(cx - d / 2), IN(cy - d / 2), IN(d), IN(d))
    o.fill.background()
    o.line.color.rgb = color
    o.line.width = Pt(weight)
    _flat(o)
    return o


def emblem(s, cx, cy, d=1.2):
    """Draw the AKKA circular seal, or embed the real logo if present."""
    if LOGO.exists():
        s.shapes.add_picture(str(LOGO), IN(cx - d / 2), IN(cy - d / 2), IN(d), IN(d))
        return
    disc = s.shapes.add_shape(MSO_SHAPE.OVAL, IN(cx - d / 2), IN(cy - d / 2), IN(d), IN(d))
    disc.fill.solid()
    disc.fill.fore_color.rgb = CREAM
    disc.line.color.rgb = LOGO_BLUE
    disc.line.width = Pt(2.5)
    _flat(disc)
    ring(s, cx, cy, d * 0.80, LOGO_BLUE, 1.5)
    text(s, IN(cx - d / 2), IN(cy - d * 0.16), IN(d), IN(d * 0.32),
         [("\u0c85\u0c95\u0ccd\u0c95", 15, LOGO_RED, True, FONT_HEAD)],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    text(s, IN(cx - d / 2), IN(cy + d * 0.10), IN(d), IN(d * 0.26),
         [("AKKA", 11, LOGO_BLUE, True, FONT_BODY)],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)


def kicker(s, label, color=SAFFRON):
    band(s, IN(0.9), IN(0.9), IN(0.09), IN(0.34), color)
    text(s, IN(1.12), IN(0.86), IN(9), IN(0.5),
         [(label, 15, color, True, FONT_BODY)])


def footer(s, n):
    band(s, IN(0.9), IN(7.02), IN(11.53), IN(0.014), INK_2)
    text(s, IN(0.9), IN(7.08), IN(7), IN(0.35),
         [("\u0c85\u0c95\u0ccd\u0c95  Akkaverse", 11, MUTED, True, FONT_BODY)])
    text(s, IN(6.0), IN(7.08), IN(4.0), IN(0.35),
         [("Built by Mithun R", 11, MUTED, False, FONT_BODY)], align=PP_ALIGN.CENTER)
    text(s, IN(10.4), IN(7.08), IN(2.03), IN(0.35),
         [(f"{n:02d} / 10", 11, SAFFRON, True, FONT_BODY)], align=PP_ALIGN.RIGHT)


# ============================================================ SLIDE 1 — TITLE
s = slide()
bg(s)
band(s, 0, 0, IN(0.26), H, SAFFRON)
band(s, IN(0.26), 0, IN(0.09), H, MAGENTA)
glow(s, 7.4, -2.0, 6.5, MAGENTA, 12)
glow(s, 9.5, 3.6, 5.5, SAFFRON, 10)
if HERO.exists():
    frame = rounded(s, IN(9.05), IN(0.7), IN(3.65), IN(6.1), INK_2, radius=0.05)
    frame.line.color.rgb = SAFFRON
    frame.line.width = Pt(1.5)
    pic = s.shapes.add_picture(str(HERO), IN(9.18), IN(0.83), height=IN(5.84))
    if pic.width > IN(3.4):
        over = (pic.width - IN(3.4)) / pic.width / 2
        pic.crop_left = over
        pic.crop_right = over
        pic.left = IN(9.18)
        pic.width = IN(3.4)
    text(s, IN(9.05), IN(6.42), IN(3.65), IN(0.4),
         [("Three generations, one voice \u2014 preserved.", 11, MUTED, False, FONT_BODY, True)],
         align=PP_ALIGN.CENTER)
emblem(s, 1.4, 1.35, 1.15)
text(s, IN(2.15), IN(0.98), IN(6), IN(0.8),
     [("\u0c85\u0c95\u0ccd\u0c95 \u00b7 AKKAVERSE", 19, TEAL, True, FONT_BODY)], anchor=MSO_ANCHOR.MIDDLE)
text(s, IN(0.85), IN(2.35), IN(8.2), IN(2.5),
     [("Preserve.\nTeach. Celebrate.", 60, CREAM, True, FONT_HEAD)], line_spacing=1.0)
text(s, IN(0.9), IN(4.95), IN(7.9), IN(1.1),
     [("A bilingual (English + \u0c95\u0ca8\u0ccd\u0ca8\u0ca1) local-first platform that keeps "
       "Kannada language, family memory and heritage alive \u2014 one voice at a time.",
       18, MUTED, False, FONT_BODY)], line_spacing=1.15)
text(s, IN(0.9), IN(6.25), IN(7.9), IN(0.5), [("by Mithun R", 17, SAFFRON, True, FONT_HEAD)])
text(s, IN(0.9), IN(6.75), IN(7.9), IN(0.5),
     [("akkaverse.vercel.app  \u00b7  github.com/Mithuncoding/akkaverse", 12.5, MUTED, False, FONT_BODY)])

# =========================================================== SLIDE 2 — PROBLEM
s = slide()
bg(s)
kicker(s, "01  \u00b7  THE PROBLEM", MAGENTA)
text(s, IN(0.9), IN(1.3), IN(11.5), IN(1.4),
     [("Culture doesn't vanish overnight.\nIt fades one generation at a time.",
       36, CREAM, True, FONT_HEAD)], line_spacing=1.05)
gen = [
    ("\u0c85\u0c9c\u0ccd\u0c9c\u0cbf", "Grandparent", "Speaks Kannada fluently", SAFFRON),
    ("Parent", "Understands", "But rarely speaks it", MAGENTA),
    ("Child", "A few words", "The thread is breaking", TEAL),
]
cx = 0.9
for hd, sub, body, col in gen:
    rounded(s, IN(cx), IN(3.15), IN(3.7), IN(2.45), CARD)
    band(s, IN(cx), IN(3.15), IN(3.7), IN(0.14), col)
    text(s, IN(cx + 0.32), IN(3.5), IN(3.1), IN(0.7), [(hd, 24, col, True, FONT_HEAD)])
    text(s, IN(cx + 0.32), IN(4.24), IN(3.1), IN(0.6), [(sub, 18, CREAM, True, FONT_BODY)])
    text(s, IN(cx + 0.32), IN(4.88), IN(3.1), IN(0.8), [(body, 14, MUTED, False, FONT_BODY)])
    cx += 3.95
text(s, IN(0.9), IN(6.0), IN(11.5), IN(0.95),
     [("When an elder is gone, their voice, blessings, proverbs, recipes and stories "
       "go too. Existing heritage apps are static encyclopedias \u2014 they store facts, "
       "not living family memory.", 16, MUTED, False, FONT_BODY)], line_spacing=1.15)
footer(s, 2)

# ========================================================== SLIDE 3 — SOLUTION
s = slide()
bg(s)
kicker(s, "02  \u00b7  THE SOLUTION", SAFFRON)
text(s, IN(0.9), IN(1.3), IN(11.5), IN(1.1),
     [("Start with your family \u2014 not a textbook.", 36, CREAM, True, FONT_HEAD)])
steps = [
    ("Preserve", "Capture an elder's words (blessing, story, song, recipe, advice) "
     "in Kannada + English, with their original recording.", SAFFRON),
    ("Teach", "Akkaverse auto-generates a bilingual micro-lesson \u2014 a phrase to "
     "repeat, key vocabulary and a meaning check for the next generation.", MAGENTA),
    ("Celebrate", "Explore the wider culture with a grounded AI guide, an interactive "
     "map, a cinematic timeline, festivals, stories and quizzes.", TEAL),
]
y = 2.55
for i, (hd, body, col) in enumerate(steps, 1):
    rounded(s, IN(0.9), IN(y), IN(11.5), IN(1.28), CARD)
    d = s.shapes.add_shape(MSO_SHAPE.OVAL, IN(1.2), IN(y + 0.32), IN(0.66), IN(0.66))
    d.fill.solid(); d.fill.fore_color.rgb = col; _no_line(d); _flat(d)
    text(s, IN(1.2), IN(y + 0.34), IN(0.66), IN(0.66), [(str(i), 24, INK, True, FONT_HEAD)],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, IN(2.15), IN(y + 0.16), IN(2.4), IN(1), [(hd, 23, col, True, FONT_HEAD)],
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, IN(4.5), IN(y + 0.12), IN(7.6), IN(1.05), [(body, 15, CREAM, False, FONT_BODY)],
         anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
    y += 1.5
footer(s, 3)

# ============================================================ SLIDE 4 — DEMO
s = slide()
bg(s)
kicker(s, "03  \u00b7  PRODUCT DEMO", TEAL)
text(s, IN(0.9), IN(1.3), IN(11.5), IN(1),
     [("Four surfaces. One heritage universe.", 36, CREAM, True, FONT_HEAD)])
cards = [
    ("Ask Akka", "AI guide with cited Wikipedia answers in EN / \u0c95\u0ca8\u0ccd\u0ca8\u0ca1.", SAFFRON),
    ("Voice Legacy", "Preserve a voice \u2192 auto bilingual lesson \u2192 pass it on.", MAGENTA),
    ("Explore Map", "Interactive Karnataka district map + photos.", TEAL),
    ("Timeline", "Cinematic scroll through 2,500 years of history.", SAFFRON),
]
positions = [(0.9, 2.55), (6.75, 2.55), (0.9, 4.72), (6.75, 4.72)]
for (hd, body, col), (px, py) in zip(cards, positions):
    rounded(s, IN(px), IN(py), IN(5.65), IN(1.9), CARD)
    band(s, IN(px), IN(py), IN(0.14), IN(1.9), col)
    text(s, IN(px + 0.4), IN(py + 0.28), IN(5), IN(0.6), [(hd, 23, col, True, FONT_HEAD)])
    text(s, IN(px + 0.4), IN(py + 0.96), IN(5), IN(0.9), [(body, 15, CREAM, False, FONT_BODY)],
         line_spacing=1.1)
text(s, IN(0.9), IN(6.78), IN(11.5), IN(0.4),
     [("Fully bilingual  \u00b7  no login  \u00b7  works offline (PWA)", 14, MUTED, True, FONT_BODY)])
footer(s, 4)

# ========================================================= SLIDE 5 — HOW IT WORKS
s = slide()
bg(s)
kicker(s, "04  \u00b7  HOW IT WORKS", SAFFRON)
text(s, IN(0.9), IN(1.3), IN(11.5), IN(1),
     [("Grounded AI \u2014 not a chatbot wrapper.", 36, CREAM, True, FONT_HEAD)])
flow = [
    ("User asks", "EN or \u0c95\u0ca8\u0ccd\u0ca8\u0ca1", MUTED),
    ("Retrieve", "Wikipedia REST\n+ PageImages", TEAL),
    ("Ground", "Inject context\n+ citations", SAFFRON),
    ("NVIDIA NIM", "Llama 3.1 8B \u2192\nNemotron 49B", MAGENTA),
    ("Cited answer", "Streamed with\nsources", TEAL),
]
x = 0.9
for i, (hd, body, col) in enumerate(flow):
    rounded(s, IN(x), IN(2.75), IN(2.15), IN(1.75), CARD)
    band(s, IN(x), IN(2.75), IN(2.15), IN(0.12), col)
    text(s, IN(x + 0.2), IN(3.0), IN(1.8), IN(0.6), [(hd, 16, col, True, FONT_BODY)])
    text(s, IN(x + 0.2), IN(3.6), IN(1.8), IN(0.9), [(body, 12.5, CREAM, False, FONT_BODY)],
         line_spacing=1.05)
    if i < len(flow) - 1:
        text(s, IN(x + 2.15), IN(3.2), IN(0.4), IN(0.8), [("\u2192", 26, SAFFRON, True, FONT_BODY)],
             align=PP_ALIGN.CENTER)
    x += 2.42
rounded(s, IN(0.9), IN(4.95), IN(11.5), IN(1.7), INK_2)
text(s, IN(1.25), IN(5.15), IN(11), IN(1.3),
     [("Server-only API key \u2014 never shipped to the browser. Local-first storage: "
       "family data in localStorage, original audio in IndexedDB. Nothing is uploaded. "
       "If the key or network is missing, curated fallbacks still answer.",
       16, CREAM, False, FONT_BODY)], line_spacing=1.2, anchor=MSO_ANCHOR.MIDDLE)
footer(s, 5)

# ============================================================ SLIDE 6 — TECH STACK
s = slide()
bg(s)
kicker(s, "05  \u00b7  TECH STACK & TRUST", TEAL)
text(s, IN(0.9), IN(1.3), IN(11.5), IN(1),
     [("Lean stack. Serious guarantees.", 36, CREAM, True, FONT_HEAD)])
stack = [
    ("Frontend", "Next.js (App Router), React, TypeScript, Tailwind, shadcn/ui"),
    ("AI / LLMs", "NVIDIA NIM \u2014 Llama 3.1 8B \u2192 Nemotron Super 49B (language-aware)"),
    ("Grounding", "Wikipedia REST + PageImages retrieval with citations"),
    ("Voice & OCR", "Keyless Kannada TTS + browser fallback; Tesseract.js in-browser"),
    ("Storage", "localStorage + IndexedDB \u2014 local-first, no login, offline PWA"),
    ("Deploy", "Vercel \u2014 single deployment, one free NIM key, no database"),
]
y = 2.6
for hd, body in stack:
    text(s, IN(0.95), IN(y), IN(2.6), IN(0.5), [(hd, 17, SAFFRON, True, FONT_BODY)])
    text(s, IN(3.7), IN(y), IN(8.6), IN(0.5), [(body, 16, CREAM, False, FONT_BODY)])
    band(s, IN(0.95), IN(y + 0.52), IN(11.4), IN(0.012), INK_2)
    y += 0.72
footer(s, 6)

# ========================================================== SLIDE 7 — DIFFERENT
s = slide()
bg(s)
kicker(s, "06  \u00b7  WHY IT'S DIFFERENT", MAGENTA)
text(s, IN(0.9), IN(1.3), IN(11.5), IN(1),
     [("Not another encyclopedia app.", 36, CREAM, True, FONT_HEAD)])
diff = [
    ("Family-first", "Begins with your elder's voice, not a textbook.", SAFFRON),
    ("Shows its sources", "Grounded AI cites Wikipedia \u2014 no invented history.", TEAL),
    ("Private by design", "No account, no upload; data stays on your device.", MAGENTA),
    ("Truly bilingual", "Every string in EN, \u0c95\u0ca8\u0ccd\u0ca8\u0ca1, or both at once.", TEAL),
    ("Works offline", "Core pages run without a network (PWA).", SAFFRON),
    ("Zero-friction demo", "No login \u2014 judges use every feature instantly.", MAGENTA),
]
positions = [(0.9, 2.6), (6.75, 2.6), (0.9, 3.9), (6.75, 3.9), (0.9, 5.2), (6.75, 5.2)]
for (hd, body, col), (px, py) in zip(diff, positions):
    rounded(s, IN(px), IN(py), IN(5.65), IN(1.1), CARD)
    dot(s, IN(px + 0.32), IN(py + 0.42), col, 0.26)
    text(s, IN(px + 0.85), IN(py + 0.18), IN(4.6), IN(0.5), [(hd, 18, col, True, FONT_BODY)])
    text(s, IN(px + 0.85), IN(py + 0.6), IN(4.6), IN(0.5), [(body, 13.5, MUTED, False, FONT_BODY)])
footer(s, 7)

# ========================================================= SLIDE 8 — VIABILITY
s = slide()
bg(s)
kicker(s, "07  \u00b7  BUSINESS VIABILITY & SCALE", SAFFRON)
text(s, IN(0.9), IN(1.3), IN(11.5), IN(1),
     [("High margin. Real Karnataka impact.", 36, CREAM, True, FONT_HEAD)])
rounded(s, IN(0.9), IN(2.55), IN(11.5), IN(1.05), INK_2)
text(s, IN(1.2), IN(2.57), IN(11), IN(1.05),
     [("Scale path:  Diaspora families  \u2192  Kannada schools & cultural associations  "
       "\u2192  other Indian languages (same engine, swap datasets)",
       16, CREAM, True, FONT_BODY)], anchor=MSO_ANCHOR.MIDDLE)
cols = [
    ("Karnataka angle", "Tourism via the district Explore layer; a discovery surface "
     "for local MSMEs & artisans through festivals and stories.", TEAL),
    ("Cost model", "Near-zero infra \u2014 one Vercel deploy + free NIM tier, no database. "
     "High margin by construction.", SAFFRON),
    ("Revenue", "School / institution licences and an optional opt-in cloud-sync tier "
     "for families across devices.", MAGENTA),
]
x = 0.9
for hd, body, col in cols:
    rounded(s, IN(x), IN(3.85), IN(3.7), IN(2.7), CARD)
    band(s, IN(x), IN(3.85), IN(3.7), IN(0.14), col)
    text(s, IN(x + 0.32), IN(4.17), IN(3.1), IN(0.6), [(hd, 18, col, True, FONT_HEAD)])
    text(s, IN(x + 0.32), IN(4.9), IN(3.1), IN(1.5), [(body, 14, CREAM, False, FONT_BODY)],
         line_spacing=1.15)
    x += 3.95
footer(s, 8)

# ============================================================ SLIDE 9 — IMPACT
s = slide()
bg(s)
kicker(s, "08  \u00b7  IMPACT & FEEDBACK", TEAL)
text(s, IN(0.9), IN(1.3), IN(11.5), IN(1),
     [("Every preserved voice becomes a lesson.", 36, CREAM, True, FONT_HEAD)])
stats = [
    ("1 min", "to preserve a voice\nand auto-build a lesson", SAFFRON),
    ("0", "logins, uploads, or\ncredentials required", MAGENTA),
    ("EN + \u0c95\u0ca8\u0ccd\u0ca8\u0ca1", "every screen,\nfully bilingual", TEAL),
]
x = 0.9
for hd, body, col in stats:
    rounded(s, IN(x), IN(2.6), IN(3.7), IN(1.9), CARD)
    text(s, IN(x + 0.3), IN(2.8), IN(3.1), IN(0.9), [(hd, 34, col, True, FONT_HEAD)])
    text(s, IN(x + 0.32), IN(3.7), IN(3.2), IN(0.8), [(body, 13.5, CREAM, False, FONT_BODY)],
         line_spacing=1.05)
    x += 3.95
rounded(s, IN(0.9), IN(4.8), IN(11.5), IN(1.75), INK_2)
text(s, IN(1.25), IN(5.0), IN(11), IN(1.4),
     [("Inclusive by design \u2014 usable by non-technical elders on any device. The "
       "preserve-and-teach loop turns a passive \u201csomeday I'll record grandpa\u201d "
       "intention into a finished, shareable bilingual lesson.", 17, CREAM, False, FONT_BODY),
      ("Add tester quotes / feedback gathered during the build phase here.",
       13, MUTED, False, FONT_BODY, True)],
     line_spacing=1.2, anchor=MSO_ANCHOR.MIDDLE, space_after=8)
footer(s, 9)

# ============================================================ SLIDE 10 — CLOSE
s = slide()
bg(s)
band(s, 0, 0, IN(0.26), H, SAFFRON)
band(s, IN(0.26), 0, IN(0.09), H, MAGENTA)
glow(s, 8.6, 3.2, 6.0, MAGENTA, 12)
kicker(s, "09  \u00b7  6-MONTH ROADMAP & THE ASK", SAFFRON)
emblem(s, 11.7, 1.25, 1.05)
text(s, IN(0.9), IN(1.3), IN(9.5), IN(1),
     [("If we had six more months\u2026", 36, CREAM, True, FONT_HEAD)])
road = [
    ("Opt-in cloud sync", "RLS-protected accounts + cross-device audio."),
    ("Voice-cloning narration", "Preserved recordings speak new text in the elder's own voice (with consent)."),
    ("Community oral-history wall", "A moderated public archive of shared voices."),
    ("School mode", "Teacher dashboards + assignments from Voice Legacy lessons."),
    ("Kannada literary retrieval", "Grounding beyond Wikipedia into curated Kannada sources."),
]
y = 2.45
for hd, body in road:
    dot(s, IN(1.0), IN(y + 0.14), SAFFRON, 0.2)
    text(s, IN(1.4), IN(y), IN(4.1), IN(0.6), [(hd, 16, CREAM, True, FONT_BODY)])
    text(s, IN(5.5), IN(y), IN(6.8), IN(0.6), [(body, 14, MUTED, False, FONT_BODY)])
    y += 0.56
rounded(s, IN(0.9), IN(5.42), IN(11.53), IN(1.46), INK_2)
text(s, IN(1.25), IN(5.58), IN(11), IN(0.5),
     [("Thank you.", 19, SAFFRON, True, FONT_HEAD),
      ("   Let's keep every family's voice alive.", 15, CREAM, False, FONT_BODY)],
     space_after=0)
text(s, IN(1.25), IN(6.34), IN(11), IN(0.5),
     [("Live  akkaverse.vercel.app      \u00b7      Code  github.com/Mithuncoding/akkaverse"
       "      \u00b7      Video  youtu.be/lXtAyZyDumM", 13, MUTED, False, FONT_BODY)])
text(s, IN(0.9), IN(7.08), IN(11.53), IN(0.35),
     [("\u0c85\u0c95\u0ccd\u0c95  Akkaverse  \u00b7  Built by Mithun R", 12, MUTED, True, FONT_BODY)],
     align=PP_ALIGN.CENTER)

# --------------------------------------------------------------------- save
out = Path(__file__).with_name("Akkaverse-Pitch-Deck.pptx")
prs.save(out)
print(f"Saved {out}  ({len(prs.slides._sldIdLst)} slides)  logo="
      f"{'real' if LOGO.exists() else 'drawn'}  hero={'yes' if HERO.exists() else 'no'}")
